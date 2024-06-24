import numpy as np
import pandas as pd
from io import BytesIO

from pptx import Presentation

from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN


import datetime as dt


dt_string = dt.datetime.now().strftime("%d/%m/%Y %H:%M:%S")

from flask import session, request


def main(form, request):
    school_year = session["school_year"]
    year = school_year + 1
    term = session["term"]

    grad_list_file = request.files[form.grad_list_file.name]
    grad_list_df = pd.read_csv(grad_list_file).fillna("")
    grad_list_df = grad_list_df.drop_duplicates(subset=["StudentID"])
    grad_list_df = grad_list_df.sort_values(by=["LastName", "FirstName"])

    # Start Presentation
    prs = Presentation()
    blank_slid_layout = prs.slide_layouts[6]
    title_and_content_slid_layout = prs.slide_layouts[1]
    prs.slide_height = Inches(8.5)
    prs.slide_width = Inches(11)

    # Iterate through students
    dt_string = dt.datetime.now().strftime("%m/%d/%Y %H:%M:%S")

    for index, student in grad_list_df.iterrows():

        slide = prs.slides.add_slide(blank_slid_layout)
        first_name = student["FirstName"].title()
        last_name = student["LastName"].title()

        june_grad = student["June Grad?"]
        is_transcript_finalized = student["Transcript Finalized?"]
        june_diploma_type = student["Diploma Type"]
        regents_with_honors = student["With Honors (Regents)?"]

        ## Update Time
        shapes = slide.shapes
        left = Inches(1.75)
        top = Inches(7.5)
        height = Inches(0.25)
        width = Inches(7.5)
        shape = shapes.add_textbox(left, top, width, height)
        line = shape.line
        line.fill.background()
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"Updated on: {dt_string}"
        font = run.font
        font.name = "Playfair Display"
        font.size = Pt(14)

        # Graduate?
        if june_grad in ["Y", ""]:
            print(is_transcript_finalized)
            if june_diploma_type == "" or is_transcript_finalized == False:
                shapes = slide.shapes
                left = Inches(1.75)
                top = Inches(0.32)
                height = Inches(1)
                width = Inches(7.5)
                shape = shapes.add_textbox(left, top, width, height)
                line = shape.line
                line.fill.background()
                text_frame = shape.text_frame
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = f"June {year} Graduate - DIPLOMA NOT FINALIZED"
                font = run.font
                font.name = "Playfair Display"
                font.size = Pt(36)
        if june_grad in ["N", "N-August"]:
            shapes = slide.shapes
            left = Inches(1.75)
            top = Inches(0.32)
            height = Inches(1)
            width = Inches(7.5)
            shape = shapes.add_textbox(left, top, width, height)
            line = shape.line
            line.fill.background()
            text_frame = shape.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = f"Not June {year} Graduate - {june_grad}"
            font = run.font
            font.name = "Playfair Display"
            font.size = Pt(36)

        ## Add Student Name
        shapes = slide.shapes
        left = Inches(1.62)
        top = Inches(3.7)
        height = Inches(1.11)
        width = Inches(7.77)
        shape = shapes.add_textbox(left, top, width, height)
        line = shape.line
        line.fill.background()
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"{first_name} {last_name}"
        font = run.font
        font.name = "Playfair Display"
        font.size = Pt(60)
        font.color.rgb = RGBColor(0x00, 0x00, 0x00)

        ## Add Diploma Type
        shapes = slide.shapes
        left = Inches(0.5)
        top = Inches(6.5)
        height = Inches(0.5)
        width = Inches(2)
        shape = shapes.add_textbox(left, top, width, height)
        line = shape.line
        line.fill.background()
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        font = run.font
        font.name = "Playfair Display"
        font.size = Pt(18)

        if june_diploma_type in ["R"]:
            run.text = f"Regents"
        elif june_diploma_type in ["AR"]:
            run.text = f"Advanced Regents"
            if regents_with_honors == True:
                shapes = slide.shapes
                left = Inches(0.5)
                top = Inches(7)
                height = Inches(0.5)
                width = Inches(2)
                shape = shapes.add_textbox(left, top, width, height)
                line = shape.line
                line.fill.background()
                text_frame = shape.text_frame
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                font = run.font
                font.name = "Playfair Display"
                font.size = Pt(18)
                run.text = f"With Honors"
        elif june_diploma_type in ["L"]:
            run.text = f"Local"

        ## Add GPA Honors
        shapes = slide.shapes
        left = Inches(3.2)
        top = Inches(6.75)
        height = Inches(0.5)
        width = Inches(0.75)
        shape = shapes.add_textbox(left, top, width, height)
        line = shape.line
        line.fill.background()
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        font = run.font
        font.name = "Playfair Display"
        font.size = Pt(18)

        with_merit = student["With Merit (GPA)?"]
        if with_merit == True:
            run.text = f"With Merit (GPA)"
        with_honors = student["With Honors (GPA)?"]
        if with_honors == True:
            run.text = f"With Honors (GPA)"

        ## Add CTE Endorsement
        cte_endorsement = student["CTE Endorsed?"]
        if cte_endorsement == True:
            shapes = slide.shapes
            left = Inches(7.75)
            top = Inches(6.5)
            height = Inches(0.5)
            width = Inches(0.75)
            shape = shapes.add_textbox(left, top, width, height)
            line = shape.line
            line.fill.background()
            text_frame = shape.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            font = run.font
            font.name = "Playfair Display"
            font.size = Pt(18)
            run.text = f"CTE"
        ## Add Math Endorsement
        math_endorsement = student["Math Endorsed?"]
        if math_endorsement == True and june_diploma_type in ["AR"]:
            shapes = slide.shapes
            left = Inches(7.75)
            top = Inches(7.25)
            height = Inches(0.5)
            width = Inches(0.75)
            shape = shapes.add_textbox(left, top, width, height)
            line = shape.line
            line.fill.background()
            text_frame = shape.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            font = run.font
            font.name = "Playfair Display"
            font.size = Pt(18)
            run.text = f"Math"
        ## Add art Endorsement
        arts_endorsement = student["Arts Endorsed?"]
        if arts_endorsement == True:
            shapes = slide.shapes
            left = Inches(9.25)
            top = Inches(6.5)
            height = Inches(0.5)
            width = Inches(0.75)
            shape = shapes.add_textbox(left, top, width, height)
            line = shape.line
            line.fill.background()
            text_frame = shape.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            font = run.font
            font.name = "Playfair Display"
            font.size = Pt(18)
            run.text = f"Art"
        ## Add Science Endorsement
        science_endorsement = student["Science Endorsed?"]
        if science_endorsement == True and june_diploma_type in ["AR"]:
            shapes = slide.shapes
            left = Inches(9.25)
            top = Inches(7.25)
            height = Inches(0.5)
            width = Inches(0.75)
            shape = shapes.add_textbox(left, top, width, height)
            line = shape.line
            line.fill.background()
            text_frame = shape.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            font = run.font
            font.name = "Playfair Display"
            font.size = Pt(18)
            run.text = f"Science"

        ## Add Seal of Civic Readiness
        endorsement_title = "Seal of Civic Readiness"
        endorsement_flag = student[endorsement_title + "?"]
        if endorsement_flag:
            shapes = slide.shapes
            left = Inches(7.25)
            top = Inches(7.75)
            height = Inches(0.5)
            width = Inches(3.5)
            shape = shapes.add_textbox(left, top, width, height)
            line = shape.line
            line.fill.background()
            text_frame = shape.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            font = run.font
            font.name = "Playfair Display"
            font.size = Pt(18)
            run.text = endorsement_title

        ## Add NHS
        endorsement_title = "NHS"
        endorsement_flag = student[endorsement_title + "?"]
        if endorsement_flag:
            shapes = slide.shapes
            left = Inches(9.25)
            top = Inches(5.75)
            height = Inches(0.5)
            width = Inches(0.75)
            shape = shapes.add_textbox(left, top, width, height)
            line = shape.line
            line.fill.background()
            text_frame = shape.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            font = run.font
            font.name = "Playfair Display"
            font.size = Pt(18)
            run.text = endorsement_title

    # Save Presentation
    f = BytesIO()
    prs.save(f)
    f.seek(0)

    return f
